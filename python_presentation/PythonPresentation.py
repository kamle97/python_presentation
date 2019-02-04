# Kamil Lesiak

import warnings
from fpdf import FPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches

WORDS_IN_LINE = 10


def to_string(list_):
    """
    Convert given list to string (spliting with spaces) without any extra characters
    :param list_: list to convert
    :return: string with the result
    """
    result = ""
    for l in list_:
        result += l + " "
    return result


class PythonPresentation(object):
    """
    class responsible for collecting slide objects and generating presentation to pdf or pptx file
    """

    def __init__(self):
        """
        Constructor that creates empty list for slides
        """
        self.slideList = []

    def add_slides(self, *slides):
        """
        Adding slides to the presentation
        :param slides: slide or slides to add
        :return:
        """
        for slide in slides:
            self.slideList.append(slide)

    def add_slides_from_list(self, slide_list):
        """
        Adding slide list to presentation
        :param slide_list: list with slides to add
        :return:
        """
        self.slideList += slide_list

    def export_presentation_to_pptx(self, filename, slides=None):
        """
        Presentation export to the pptx format
        :param filename: name of generated pptx
        :param slides: list where you can pick slides with True values
        :return:
        """
        prs = Presentation()
        if slides is None:
            slides = [True for _ in range(len(self.slideList))]
        len_ = min(len(self.slideList), len(slides))
        if len_ != len(self.slideList):
            warnings.warn("Chosen slides list and slides list lengths are not equal!")
        for i in range(len_):
            if slides[i]:
                layout_number = self.define_layout(self.slideList[i])
                self.export_slide_content_to_pptx(prs, self.slideList[i], layout_number)
        prs.save(filename)

    def export_slide_content_to_pptx(self, prs, slide, layout):
        """
        Export single slide to pptx, function check what elements are inside the slide
        and then shows them inside pptx file
        :param prs: Presentation object inside of which slides are added
        :param slide: Slide object that determines content of a slide
        :param layout: Number of layout adopted for given slide
        :return:
        """
        slide_layout = prs.slide_layouts[layout]
        pptx_slide = prs.slides.add_slide(slide_layout)
        if slide.title is not None:
            self.set_text_style(pptx_slide.shapes.title, slide.title)
        if slide.subtitle is not None:
            self.set_text_style(pptx_slide.placeholders[1], slide.subtitle)
        if slide.text is not None:
            left = Inches(2)
            top = width = height = Inches(3)
            text_box = pptx_slide.shapes.add_textbox(left, top, width, height)
            self.set_text_style(text_box, slide.text)
        if slide.footer is not None:
            left = Inches(1)
            width = height = Inches(3)
            top = Inches(6)
            text_box = pptx_slide.shapes.add_textbox(left, top, width, height)
            self.set_text_style(text_box, slide.footer)
        if slide.image is not None:
            left = Inches(2)
            top = Inches(3)
            pptx_slide.shapes.add_picture(slide.image.filename, left, top, width=Pt(slide.image.width * 2),
                                          height=Pt(slide.image.height * 2))

    @staticmethod
    def set_text_style(shape, element):
        """
        Function editing the font properties of given shape
        :param shape: element of pptx slide that will be edited
        :param element: presentation element that holds information about the expected font look
        :return:
        """
        text_frame = shape.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = "\n".join(to_string(element.content.split()[i:i + WORDS_IN_LINE]) for i in
                             range(0, len(element.content.split()), WORDS_IN_LINE))
        font = run.font
        font.name = element.font
        font.size = Pt(element.text_size)
        font.bold = element.bold
        font.italics = element.italics
        font.underline = element.underline
        if element.text_color is not None:
            font.color.rgb = RGBColor(element.text_color.r, element.text_color.g, element.text_color.b)

    @staticmethod
    def define_layout(slide):
        """
        basis of the slide content returns the type of layout
        :param slide: slide with content
        :return: number of layout for the slide (0 - only title and subtitle, 1 - title on top, content below)
        """
        if slide.text is None and slide.image is None:
            return 0
        else:
            return 1

    def export_presentation_to_pdf(self, filename, slides=None):
        """
        Function exporting the presentation to the pdf format
        :param filename: expected pdf filename
        :param slides: list where you can pick slides with True values
        :return:
        """
        pdf = FPDF()
        if slides is None:
            slides = [True for _ in range(len(self.slideList))]
        len_ = min(len(self.slideList), len(slides))
        if len_ != len(self.slideList):
            warnings.warn("Chosen slides list and slides list lengths are not equal!")
        for i in range(len_):
            if slides[i]:
                pdf.add_page(orientation="horizontal")
                self.export_slide_content_to_pdf(pdf, self.slideList[i])

        pdf.output(filename)

    def export_slide_content_to_pdf(self, pdf, slide):
        """
        Export single slide to pdf, function check what elements are inside the slide
        and then shows them inside pdf file
        :param pdf: FPDF object, determines the file function work on
        :param slide: Slide object determines the content of a slide
        :return:
        """
        if slide.title is not None:
            self.export_to_pdf_element_with_text(pdf, slide.title)
        if slide.subtitle is not None:
            self.export_to_pdf_element_with_text(pdf, slide.subtitle)
        if slide.footer is not None:
            self.export_to_pdf_element_with_text(pdf, slide.footer)
        if slide.text is not None:
            self.export_to_pdf_element_with_text(pdf, slide.text)
        if slide.image is not None:
            pdf.image(slide.image.filename, x=slide.image.x, y=slide.image.y, w=slide.image.width, h=slide.image.height)

    def export_to_pdf_element_with_text(self, pdf, element):
        """
        Function setting the style of text
        :param pdf: presentation the function is working on
        :param element: slide element that contains style of the text
        :return:
        """
        style_text = self.read_style(element)
        pdf.set_font(element.font, style_text, element.text_size)
        self.set_text_color(pdf, element)
        pdf.set_y(element.vertical_value)
        pdf.multi_cell(0, 20, element.content, align=element.horizontal)

    @staticmethod
    def read_style(content):
        """
        Function converting boolean values of Bold/Italics/Underline to expected characters
        :param content: element contains boolean information
        :return: string with expected letters
        """
        result = ""
        if content.bold:
            result += "B"
        if content.italics:
            result += "I"
        if content.underline:
            result += "U"
        return result

    @staticmethod
    def set_text_color(pdf, element):
        """
        function setting the color of the text, default set it to black
        :param pdf: file the function is working on
        :param element: element contains color RGB information
        :return:
        """
        if element.text_color is not None:
            pdf.set_text_color(element.text_color.r, element.text_color.g, element.text_color.b)
        else:
            pdf.set_text_color(0, 0, 0)
